# -*- coding: utf-8 -*-
"""Genera el mazo EXPO-E3_ACC_Carpio-Pacheco-Cando-Alvarez.pptx segun la rubrica
EX-PFC-E3-DIAP. 18 diapositivas totales (15 de contenido + portada + declaracion IA +
referencias), 16:9, Arial, paleta de 3 colores semanticos + neutros con contraste WCAG
verificado, regla de 2 lineas, numeracion n/N, pie fijo, marcadores de seccion,
rotulacion de expositor por bloque.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# Rutas relativas al repo (este script vive en docs/diapositivas/) -- para que el mazo
# se pueda regenerar en cualquier clon del repositorio, no solo en esta maquina.
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(SCRIPT_DIR, "..", ".."))

def repo_path(*parts):
    return os.path.join(REPO_ROOT, *parts)

# ---------------------------------------------------------------- Paleta / constantes
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
TEAL = RGBColor(0x0F, 0x6B, 0x5C)
AMBER = RGBColor(0x8A, 0x5A, 0x00)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x59, 0x59, 0x59)
FOOTER_BG = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BG = RGBColor(0xF2, 0xF4, 0xF7)
FONT = "Arial"

EQUIPO = "Equipo ACC"
SISTEMA = "Sistema de Soporte Técnico ISP"
CODIGO = "PFC-E3 · ISR-701"
REPO_URL = "https://github.com/carlospatroner-boop/PFC-SOPORTE-ISP"

# Geometria (16:9, 33.87 x 19.05 cm = 13.333 x 7.5 in)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(1 / 2.54)  # 1.0 cm de seguridad
USABLE_L = MARGIN
USABLE_T = MARGIN
USABLE_W = SLIDE_W - 2 * MARGIN
USABLE_R = SLIDE_W - MARGIN
USABLE_B = SLIDE_H - MARGIN

TAG_ROW_H = Inches(0.38)
TITLE_ROW_H = Inches(0.78)
BODY_TOP = USABLE_T + TAG_ROW_H + TITLE_ROW_H + Inches(0.08)
BODY_BOTTOM = Inches(7.00)
BODY_H = BODY_BOTTOM - BODY_TOP

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL_CONTENT_SLIDES = 17  # todas menos portada, para numeracion n/N (portada no numera)


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_font(run, size, color=TEXT, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def add_rect(slide, l, t, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_textbox(slide, l, t, w, h, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def top_accent_bar(slide, color=NAVY):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), color)


def section_and_presenter_tag(slide, section_label, presenter=None, tag_color=NAVY):
    # Etiqueta de seccion (marcador de avance de la estructura)
    tag_w = Inches(3.4)
    tag = add_rect(slide, USABLE_L, USABLE_T, tag_w, TAG_ROW_H, tag_color)
    tf = tag.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.02)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section_label
    set_font(r, 14, WHITE, bold=True)
    # Rotulacion del expositor responsable del bloque
    if presenter:
        _, tf2 = add_textbox(slide, USABLE_R - Inches(3.6), USABLE_T, Inches(3.6), TAG_ROW_H,
                              valign=MSO_ANCHOR.MIDDLE)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = f"Expone: {presenter}"
        set_font(r2, 14, GRAY, italic=True)


def title(slide, text, color=NAVY, size=32):
    top = USABLE_T + TAG_ROW_H + Inches(0.06)
    _, tf = add_textbox(slide, USABLE_L, top, USABLE_W, TITLE_ROW_H, valign=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_font(r, size, color, bold=True)


def footer(slide, page_num=None, show_number=True):
    bar_t = Inches(7.14)
    bar_h = SLIDE_H - bar_t
    add_rect(slide, 0, bar_t, SLIDE_W, bar_h, FOOTER_BG)
    _, tf = add_textbox(slide, Inches(0.3), bar_t, Inches(9.5), bar_h, valign=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{EQUIPO} · {CODIGO} · {SISTEMA}"
    set_font(r, 14, WHITE)
    if show_number and page_num is not None:
        _, tf2 = add_textbox(slide, SLIDE_W - Inches(1.6), bar_t, Inches(1.3), bar_h,
                              valign=MSO_ANCHOR.MIDDLE)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = f"{page_num}/{TOTAL_CONTENT_SLIDES}"
        set_font(r2, 14, WHITE, bold=True)


def bullet_panel(slide, bullets, top=BODY_TOP, height=BODY_H, accent=NAVY, font_size=22):
    """Panel con fondo compositivo (no texto flotando en blanco) + vinetas de <=2 lineas."""
    panel = add_rect(slide, USABLE_L, top, USABLE_W, height, PANEL_BG)
    tf = panel.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2
        p.space_after = Pt(16)
        r = p.add_run()
        r.text = f"●  {b}"
        set_font(r, font_size, TEXT)
        # marcador de color en el bullet (icono textual ya incluido); resaltar la vineta
        r.font.color.rgb = accent if i == 0 else TEXT
    return panel


def image_protagonist(slide, path, caption=None, top=BODY_TOP, height=BODY_H, accent=NAVY):
    from PIL import Image as PILImage
    im = PILImage.open(path)
    iw, ih = im.size
    aspect = iw / ih
    avail_h = height - (Inches(0.35) if caption else 0)
    avail_w = USABLE_W
    # ajustar a la caja disponible manteniendo proporcion, maximizando area (>=50% exigido)
    h_if_full_w = avail_w / aspect
    if h_if_full_w <= avail_h:
        w, h = avail_w, h_if_full_w
    else:
        w, h = avail_h * aspect, avail_h
    left = USABLE_L + (USABLE_W - w) / 2
    slide.shapes.add_picture(path, int(left), int(top), width=int(w), height=int(h))
    if caption:
        _, tf = add_textbox(slide, USABLE_L, top + h + Inches(0.06), USABLE_W, Inches(0.32))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        set_font(r, 15, GRAY, italic=True)


def citation(slide, text, top=None):
    top = top or Inches(6.68)
    _, tf = add_textbox(slide, USABLE_L, top, USABLE_W, Inches(0.28))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_font(r, 13, GRAY)


PAGE = {"n": 0}


def new_content_slide(section_label, ttl, presenter=None, tag_color=NAVY):
    PAGE["n"] += 1
    s = add_slide()
    top_accent_bar(s, tag_color)
    section_and_presenter_tag(s, section_label, presenter, tag_color)
    title(s, ttl)
    footer(s, PAGE["n"])
    return s


# ============================================================ 1) PORTADA
s = add_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
_, tf = add_textbox(s, Inches(1.0), Inches(1.7), Inches(11.33), Inches(1.4))
p = tf.paragraphs[0]
r = p.add_run(); r.text = SISTEMA
set_font(r, 40, WHITE, bold=True)
_, tf = add_textbox(s, Inches(1.0), Inches(2.75), Inches(11.33), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Entrega 3 · Aplicaciones Distribuidas (ISR-701) · Equipo ACC"
set_font(r, 22, RGBColor(0xCF, 0xDA, 0xE8))

integrantes = [
    "Carlos Carpio — Datos/Reportes · expone la demo",
    "Cristhian Pacheco — Ticket-service / Spark",
    "Robinson Cando — Pruebas / Métricas / Protocolo",
    "Jeremy Álvarez — Arquitectura / Notificaciones / IA",
]
_, tf = add_textbox(s, Inches(1.0), Inches(3.75), Inches(11.33), Inches(2.0))
for i, line in enumerate(integrantes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = line
    set_font(r, 18, WHITE)

add_rect(s, Inches(1.0), Inches(6.15), Inches(11.33), Inches(0.02), RGBColor(0x4A, 0x63, 0x86))
_, tf = add_textbox(s, Inches(1.0), Inches(6.35), Inches(11.33), Inches(0.4))
p = tf.paragraphs[0]
r = p.add_run(); r.text = f"Código: {CODIGO}    ·    Repositorio: {REPO_URL}"
set_font(r, 15, RGBColor(0xCF, 0xDA, 0xE8))
_, tf = add_textbox(s, Inches(1.0), Inches(6.75), Inches(11.33), Inches(0.35))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Julio 2026"
set_font(r, 14, RGBColor(0x9A, 0xAD, 0xC7))

# ============================================================ 2) HOJA DE RUTA
s = new_content_slide("SECCIÓN 1 · CONTEXTO", "Objetivo de hoy y hoja de ruta", "Carlos Carpio")
panel = add_rect(s, USABLE_L, BODY_TOP, USABLE_W, Inches(0.75), NAVY)
tf = panel.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.4)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Objetivo: demostrar la capa de datos distribuida y el pipeline paralelo de la E3."
set_font(r, 19, WHITE, bold=True)
bullet_panel(s, [
    "Arquitectura vigente y qué cambió en E3",
    "Decisiones técnicas: fragmentación, CAP/PACELC",
    "Evidencia real: pruebas, métricas y pipeline Spark",
    "Demo en vivo: crear ticket → notificación",
    "Riesgos, plan hacia E4 y conclusiones",
], top=BODY_TOP + Inches(0.9), height=BODY_H - Inches(0.9))

# ============================================================ 3) PROBLEMA Y ALCANCE
s = new_content_slide("SECCIÓN 1 · CONTEXTO", "El problema: soporte técnico de un ISP", "Cristhian Pacheco")
bullet_panel(s, [
    "Dominio: gestión de tickets de soporte técnico de Internet",
    "Usuarios: cliente, técnico de zona, administrador",
    "RNF crítico: SLA de resolución por prioridad",
    "RNF crítico: consistencia bajo escrituras concurrentes",
    "RNF crítico: escalar a más de 500k incidencias históricas",
])

# ============================================================ 4) ARQUITECTURA 1/2 (diagrama)
s = new_content_slide("SECCIÓN 2 · ARQUITECTURA", "Arquitectura de contenedores (C4 Nivel 2)", "Jeremy Álvarez")
image_protagonist(
    s, repo_path("docs", "diagrams", "c4-nivel2-diapositiva.png"),
    caption="Diagrama propio, notación C4 · ámbar = fragmentado en E3 · teal = nuevo en E3",
)

# ============================================================ 5) ARQUITECTURA 2/2
s = new_content_slide("SECCIÓN 2 · ARQUITECTURA", "Qué cambió respecto a E2", "Jeremy Álvarez")
bullet_panel(s, [
    "tickets: fragmentado por fecha_apertura (antes: por zona)",
    "Clúster CockroachDB real de 3 nodos (antes: mono-nodo)",
    "+3 métricas Prometheus incrementales en ticket-service",
    "+Testcontainers: pruebas contra CockroachDB real",
    "+Pipeline Spark: nuevo componente analítico paralelo",
])

# ============================================================ 6) ALCANCE E3 Y ESTADO REAL
s = new_content_slide("SECCIÓN 3 · AVANCE", "Lo planificado vs. lo construido", "Robinson Cando", tag_color=AMBER)
items = [
    ("✓", TEAL, "Fragmentación por fecha + ADR-0003 — hecho, verificado"),
    ("✓", TEAL, "Testcontainers + 3 métricas Prometheus — hecho, 0 advertencias"),
    ("✓", TEAL, "Pipeline Spark, 5 transformaciones, 600k filas — hecho"),
    ("⧗", AMBER, "Video de tolerancia a fallos (1 y 2 nodos) — pendiente"),
    ("⧗", AMBER, "Protocolo estadístico, 50 corridas, IC 95% — pendiente"),
]
panel = add_rect(s, USABLE_L, BODY_TOP, USABLE_W, BODY_H, PANEL_BG)
tf = panel.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
for i, (icon, color, text) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(16)
    r1 = p.add_run(); r1.text = f"{icon}  "
    set_font(r1, 22, color, bold=True)
    r2 = p.add_run(); r2.text = text
    set_font(r2, 22, TEXT)

# ============================================================ 7) DECISIONES 1/2 (codigo)
s = new_content_slide("SECCIÓN 4 · DECISIONES", "Decisión: fragmentación por fecha_apertura", "Cristhian Pacheco", tag_color=TEAL)
code_lines = [
    "CREATE TABLE tickets (",
    "    created_at TIMESTAMPTZ NOT NULL DEFAULT now(),",
    "    id UUID NOT NULL DEFAULT gen_random_uuid(),",
    "    ...",
    "    PRIMARY KEY (created_at, id)",
    ") PARTITION BY RANGE (created_at) (",
    "    PARTITION tickets_2026_q1",
    "      VALUES FROM (MINVALUE) TO ('2026-04-01'),",
    "    ...",
    ");",
]
code_box = add_rect(s, USABLE_L, BODY_TOP, USABLE_W, Inches(3.05), RGBColor(0x1E, 0x1E, 0x22))
tf = code_box.text_frame
tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.2)
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.15
    r = p.add_run(); r.text = line
    highlight = "PARTITION BY RANGE" in line or "tickets_2026_q1" in line
    set_font(r, 16, AMBER if highlight else RGBColor(0xE5, 0xE5, 0xE5), bold=highlight, name="Consolas")
_, tf2 = add_textbox(s, USABLE_L, BODY_TOP + Inches(3.25), USABLE_W, Inches(1.15))
p = tf2.paragraphs[0]
r = p.add_run()
r.text = "Alternativa descartada: PARTITION BY LIST(zona)."
set_font(r, 19, TEXT, bold=True)
p2 = tf2.add_paragraph()
r2 = p2.add_run()
r2.text = "Criterio: exigencia explícita de la guía de E3 para el equipo ACC [1]."
set_font(r2, 19, GRAY)

# ============================================================ 8) DECISIONES 2/2 (CAP/PACELC)
s = new_content_slide("SECCIÓN 4 · DECISIONES", "Ubicación CAP / PACELC del sistema", "Cristhian Pacheco", tag_color=TEAL)
bullet_panel(s, [
    "CockroachDB = sistema CP: prioriza consistencia [2]",
    "Ante partición: rechaza escrituras conflictivas, no arriesga",
    "Consenso Raft: quórum 2 de 3, tolera 1 nodo caído [4]",
    "Verificado en vivo: conflicto real → SQLSTATE 40001",
])

# ============================================================ 9) EVIDENCIA DE IMPLEMENTACIÓN
s = new_content_slide("SECCIÓN 5 · EVIDENCIA", "Artefactos verificables en el repositorio", "Carlos Carpio")
bullet_panel(s, [
    "db-cluster/: 3 nodos, esquema fragmentado, ADR-0003",
    "services/svc-principal: refactor completo + Testcontainers",
    "spark/: pipeline + baseline + protocolo (600.000 filas)",
    "tests/integration: 2 pruebas contra CockroachDB real",
])

# ============================================================ 9b) EVIDENCIA EN VIVO: METRICAS
s = new_content_slide("SECCIÓN 5 · EVIDENCIA", "Evidencia en vivo: métricas Prometheus reales", "Robinson Cando")
code_lines2 = [
    "$ curl localhost:8002/actuator/prometheus",
    "",
    "# HELP crdb_pool_active_connections Conexiones activas del pool",
    "# TYPE crdb_pool_active_connections gauge",
    'crdb_pool_active_connections{application="ticket-service"} 0.0',
    "",
    "# HELP crdb_transaction_retries_total Reintentos por conflicto",
    "# TYPE crdb_transaction_retries_total counter",
    'crdb_transaction_retries_total{application="ticket-service"} 0.0',
    "",
    "$ promtool check metrics < salida.prom   →   0 advertencias",
]
code_box = add_rect(s, USABLE_L, BODY_TOP, USABLE_W, BODY_H, RGBColor(0x1E, 0x1E, 0x22))
tf = code_box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.15)
for i, line in enumerate(code_lines2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.2
    r = p.add_run(); r.text = line
    is_cmd = line.startswith("$")
    is_result = "0 advertencias" in line
    color = TEAL if is_cmd else (AMBER if is_result else RGBColor(0xE5, 0xE5, 0xE5))
    set_font(r, 17, color, bold=(is_cmd or is_result), name="Consolas")

# ============================================================ 10) RESULTADOS 1/2 (verificacion)
s = new_content_slide("SECCIÓN 6 · RESULTADOS", "Verificación de corrección (medida, no supuesta)", "Robinson Cando", tag_color=AMBER)
items = [
    "Aislamiento serializable: conflicto real → SQLSTATE 40001",
    "3 métricas Prometheus: 0 advertencias de promtool",
    "Testcontainers: 2 de 2 pruebas en verde, CockroachDB real",
]
panel = add_rect(s, USABLE_L, BODY_TOP, USABLE_W, BODY_H, PANEL_BG)
tf = panel.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
for i, text in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(18)
    r1 = p.add_run(); r1.text = "✓  "
    set_font(r1, 24, TEAL, bold=True)
    r2 = p.add_run(); r2.text = text
    set_font(r2, 24, TEXT)

# ============================================================ 11) RESULTADOS 2/2 (spark chart)
s = new_content_slide("SECCIÓN 6 · RESULTADOS", "Pipeline Spark sobre 600.000 incidencias", "Cristhian Pacheco", tag_color=TEAL)
image_protagonist(
    s, repo_path("spark", "results", "resultados_pipeline_diapositiva.png"),
    caption="Ejecución única, determinista (seed=42) · spark_pipeline.ipynb, 5 transformaciones [5]",
)

# ============================================================ 12) ANCLA DE LA DEMOSTRACIÓN
s = new_content_slide("SECCIÓN 7 · DEMO EN VIVO", "Demo: crear ticket → Kafka → notificación", "Carlos Carpio")
bullet_panel(s, [
    "Objetivo: mostrar la integración asíncrona real entre servicios",
    "Paso 1: POST /tickets con JWT de un cliente autenticado",
    "Paso 2: ticket-service publica ticket.created en Kafka",
    "Paso 3: notification-service consume y genera la notificación",
    "Resultado esperado: notificación visible en menos de 2 segundos",
])

# ============================================================ 13) LECTURA DEL RESULTADO (+ contingencia)
s = new_content_slide("SECCIÓN 7 · DEMO EN VIVO", "Resultado observado (respaldo de contingencia)", "Carlos Carpio")
image_protagonist(
    s, repo_path("docs", "evidencias", "captura_demo_ticket_notificacion.png"),
    caption="Ejecutado en vivo el 26/07/2026 · ticket 2129d338-... · valida el objetivo del Paso 12",
)

# ============================================================ 14) RIESGOS / DEUDA / PLAN E4
s = new_content_slide("SECCIÓN 8 · CIERRE", "Riesgos, deuda técnica y plan hacia E4", "Robinson Cando", tag_color=AMBER)
bullet_panel(s, [
    "Riesgo: zone cruza 4 particiones — mitigado con índice secundario",
    "Deuda técnica: sin API Gateway único — planificado para E4",
    "Pendiente: video de tolerancia a fallos + protocolo estadístico",
    "E4 reutiliza sin modificar: clúster, 3 métricas, pipeline Spark",
])

# ============================================================ 15) CONCLUSIONES PROPIAS
s = new_content_slide("SECCIÓN 8 · CIERRE", "Conclusiones del equipo", "Jeremy Álvarez", tag_color=AMBER)
bullet_panel(s, [
    "Fragmentación exigida ≠ patrón de acceso real: documentar el costo, no forzarlo",
    "Verificación empírica (Testcontainers, promtool) halló un problema real antes que un usuario",
    "Pipeline con objetivo propio del dominio > transformaciones genéricas",
], font_size=22)

# ============================================================ 16) DECLARACIÓN DE USO DE IA
s = new_content_slide("SECCIÓN 8 · CIERRE", "Declaración de uso de IA generativa", tag_color=GRAY)
bullet_panel(s, [
    "Herramienta: Claude Code (Anthropic), asistente conversacional",
    "Propósito: scaffolding, implementación guiada, depuración y pruebas",
    "Alcance: todas las diapositivas — contenido revisado y verificado",
    "Decisiones de arquitectura y verificación funcional: siempre del equipo",
], font_size=20)

# ============================================================ 17) REFERENCIAS
s = new_content_slide("SECCIÓN 8 · CIERRE", "Referencias", tag_color=GRAY)
refs = [
    "[1] M. T. Özsu and P. Valduriez, Principles of Distributed Database Systems, 4th ed. Springer, 2020.",
    "[2] E. Brewer, \"CAP twelve years later,\" IEEE Computer, vol. 45, no. 2, pp. 23–29, 2012.",
    "[3] D. J. Abadi, \"Consistency tradeoffs...,\" IEEE Computer, vol. 45, no. 2, pp. 37–42, 2012.",
    "[4] D. Ongaro and J. Ousterhout, \"In search of an understandable consensus algorithm,\" USENIX ATC, 2014.",
    "[5] M. Zaharia et al., \"Resilient Distributed Datasets,\" NSDI, 2012.",
]
panel = add_rect(s, USABLE_L, BODY_TOP, USABLE_W, BODY_H, PANEL_BG)
tf = panel.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.45)
for i, ref in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = ref
    set_font(r, 16, TEXT)

out_dir = repo_path("docs", "diapositivas")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "EXPO-E3_ACC_Carpio-Pacheco-Cando-Alvarez.pptx")
prs.save(out_path)
print("saved", out_path, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
